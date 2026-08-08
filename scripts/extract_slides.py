#!/usr/bin/env python3
"""
gao-wepost-ppt-skill 幻灯片提取器

输入: 1-N 份 .ppt/.pptx/.pdf
流程:
  1. PPTX/PPT → PDF (优先 Microsoft PowerPoint，字体/颜色保真更好；
     缺失或失败时用 LibreOffice soffice 兜底)
  2. PDF 逐页渲染 → PNG (pymupdf, 自动安装)
  3. 视频页过滤: 黑像素占比 > --dark-ratio 且 文本为空/仅时间码 → 过滤
  4. 文本提取 (pymupdf; pptx 无转换器时用 python-pptx 兜底)
输出: {outdir}/{stem}/slides/slide-{NN}.png + manifest.json

用法:
  python3 extract_slides.py -i deck.pdf [-i deck2.pptx ...] -o /tmp/ppt-extract/ \
    [--dark-ratio 0.80] [--brightness-threshold 30] [--dpi 150]
"""

import argparse
import json
import os
import re
import shutil
import subprocess
import sys

try:
    from PIL import Image
except ImportError:
    sys.exit("缺少 Pillow，请先运行: python3 -m pip install --user Pillow")

DEFAULT_THRESHOLD = 30   # 灰度亮度阈值
DEFAULT_RATIO = 0.80     # 黑像素占比阈值
DEFAULT_DPI = 150
SAMPLE_WIDTH = 160       # 黑占比检测的采样宽度（加速）
TIMECODE_RE = re.compile(r"^(?:[0-9]{1,2}:){1,2}[0-9]{1,2}$")
VIDEO_KEYWORDS = ("播放", "暂停", "视频", "video", "play", "pause", "▶", "⏸", "►")


def ensure_pymupdf():
    """确保 pymupdf 可用，缺失时尝试自动安装。"""
    try:
        import fitz  # noqa: F401
        return
    except ImportError:
        pass
    print("[extract] 未找到 pymupdf，尝试自动安装...")
    r = subprocess.run(
        [sys.executable, "-m", "pip", "install", "--user", "--quiet", "pymupdf"],
        capture_output=True, text=True,
    )
    if r.returncode != 0:
        sys.exit(
            "[extract] pymupdf 安装失败。请手动运行:\n"
            "  python3 -m pip install --user pymupdf\n"
            "或安装 poppler 后用 pdftoppm: brew install poppler"
        )


def find_soffice():
    for name in ("soffice", "libreoffice"):
        p = shutil.which(name)
        if p:
            return p
    mac_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if os.path.exists(mac_path):
        return mac_path
    return None


def find_powerpoint():
    """定位 Microsoft PowerPoint（macOS 优先，字体/颜色保真优于 LibreOffice）。"""
    mac_path = "/Applications/Microsoft PowerPoint.app"
    if sys.platform == "darwin" and os.path.isdir(mac_path):
        return mac_path
    return None


def pptx_to_pdf_via_powerpoint(src, outdir):
    """用 Microsoft PowerPoint 的 AppleScript 导出 PDF（保留字体/颜色）。

    返回 PDF 路径；PowerPoint 缺失、超时或导出失败时抛 RuntimeError。
    """
    app_path = find_powerpoint()
    if not app_path:
        raise RuntimeError("未找到 Microsoft PowerPoint")
    stem = os.path.splitext(os.path.basename(src))[0]
    pdf_path = os.path.join(outdir, stem + ".pdf")
    if os.path.exists(pdf_path):
        os.remove(pdf_path)
    script = (
        'tell application "Microsoft PowerPoint"\n'
        f'  open POSIX file {json.dumps(os.path.abspath(src))}\n'
        f'  save active presentation in POSIX file {json.dumps(pdf_path)} '
        "as save as PDF\n"
        '  close active presentation saving no\n'
        "end tell\n"
    )
    try:
        r = subprocess.run(
            ["osascript", "-e", script], capture_output=True, text=True, timeout=300,
        )
    except subprocess.TimeoutExpired:
        raise RuntimeError("PowerPoint 导出超时（首次启动可能较慢）")
    if r.returncode != 0 or not os.path.exists(pdf_path):
        detail = (r.stderr or r.stdout or "").strip().splitlines()
        raise RuntimeError(
            f"PowerPoint 导出失败: {detail[-1] if detail else '未生成 PDF'}"
        )
    return pdf_path


def pptx_to_pdf(src, outdir):
    """PPTX/PPT → PDF。优先 Microsoft Office（字体/颜色保真），缺失则用 LibreOffice。"""
    if find_powerpoint():
        print(f"[extract] 使用 Microsoft PowerPoint 导出 PDF（字体/颜色保真优先）...")
        try:
            return pptx_to_pdf_via_powerpoint(src, outdir)
        except RuntimeError as e:
            print(f"[extract]  {e}，降级 LibreOffice soffice...")
    soffice = find_soffice()
    if not soffice:
        raise RuntimeError(
            f"{os.path.basename(src)} 是 PPT 格式，需要转换器。\n"
            "优先安装 Microsoft Office（字体/颜色保真最好）；\n"
            "或安装 LibreOffice: brew install --cask libreoffice\n"
            "或手动用「文件 → 导出为 PDF」转换后再输入 PDF。"
        )
    subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", outdir, src],
        check=True, capture_output=True, text=True,
    )
    pdf_path = os.path.join(outdir, os.path.splitext(os.path.basename(src))[0] + ".pdf")
    if not os.path.exists(pdf_path):
        raise RuntimeError(f"soffice 转换失败，未生成 {pdf_path}")
    return pdf_path


def pptx_texts(src):
    """python-pptx 提取每页文本（soffice 缺失时兜底）。返回 [slide_text, ...]。"""
    try:
        from pptx import Presentation
    except ImportError:
        print("[extract] 无 python-pptx，无法提取 PPT 文本（非阻塞）")
        return None
    prs = Presentation(src)
    texts = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    parts.append(t)
        texts.append("\n".join(parts))
    return texts


def pdf_texts(fitz, pdf_path):
    """pymupdf 逐页提取文本。"""
    doc = fitz.open(pdf_path)
    return [page.get_text() for page in doc]


def render_pdf(fitz, pdf_path, outdir, dpi):
    """逐页渲染 PNG，返回 [png_path, ...]（页序）。"""
    doc = fitz.open(pdf_path)
    paths = []
    for i, page in enumerate(doc, start=1):
        pix = page.get_pixmap(dpi=dpi)
        p = os.path.join(outdir, f"page-{i:03d}.png")
        pix.save(p)
        paths.append(p)
    doc.close()
    return paths


def dark_ratio(png_path, threshold=DEFAULT_THRESHOLD):
    """灰度图亮度 < threshold 的像素占比。采样缩放加速。"""
    img = Image.open(png_path).convert("L")
    w, h = img.size
    if w > SAMPLE_WIDTH:
        img = img.resize((SAMPLE_WIDTH, int(h * SAMPLE_WIDTH / w)))
    pixels = img.getdata()
    total = len(pixels)
    dark = sum(1 for p in pixels if p < threshold)
    return dark / total if total else 0.0


def is_video_text(text):
    """贝叶斯证据2: 文本为空、仅时间码、或仅播放关键词 → 判为视频页文本。"""
    if not text or not text.strip():
        return True
    tokens = [t.strip() for t in text.splitlines() if t.strip()]
    if not tokens:
        return True
    for t in tokens:
        t2 = re.sub(r"[▶⏸►▷◇○·\s:]", "", t).strip()
        if not t2:
            continue
        if TIMECODE_RE.match(t):
            continue
        if t2.lower() in VIDEO_KEYWORDS:
            continue
        return False
    return True


def process_file(src, outdir, dpi, ratio=DEFAULT_RATIO, threshold=DEFAULT_THRESHOLD):
    """处理单个输入文件，返回 manifest dict。"""
    stem = os.path.splitext(os.path.basename(src))[0]
    work = os.path.join(outdir, stem)
    tmpdir = os.path.join(work, "_tmp")
    slidir = os.path.join(work, "slides")
    os.makedirs(slidir, exist_ok=True)
    os.makedirs(tmpdir, exist_ok=True)

    ensure_pymupdf()
    import fitz

    ext = os.path.splitext(src)[1].lower()
    pdf_path = src
    texts = None
    if ext in (".ppt", ".pptx"):
        print(f"[extract] {os.path.basename(src)}: PPT → PDF（Office/LibreOffice）...")
        pdf_path = pptx_to_pdf(src, tmpdir)
        if ext == ".pptx":
            texts = pptx_texts(src)

    print(f"[extract] {os.path.basename(src)}: 渲染 {dpi}dpi...")
    pngs = render_pdf(fitz, pdf_path, tmpdir, dpi)
    if texts is None:
        texts = pdf_texts(fitz, pdf_path)
    if pdf_path != src:
        os.remove(pdf_path)

    if len(texts) != len(pngs):
        texts = [""] * len(pngs)

    slides = []
    excluded = []
    for i, (png, text) in enumerate(zip(pngs, texts), start=1):
        r = dark_ratio(png, threshold)
        video = r > ratio and is_video_text(text)
        if video:
            excluded.append(
                {"original_page": i, "dark_ratio": round(r, 3),
                 "reason": "dark_ratio>0.80 且文本为空/仅时间码"}
            )
            print(f"[extract]  page {i}: 视频页，过滤 (dark={r:.1%})")
            continue
        new_name = f"slide-{len(slides) + 1:02d}.png"
        os.rename(png, os.path.join(slidir, new_name))
        slides.append({
            "index": len(slides) + 1,
            "original_page": i,
            "file": new_name,
            "dark_ratio": round(r, 3),
            "is_video": False,
            "text": text.strip(),
        })
        flags = []
        if r > ratio:
            flags.append("深色页(黑占比超阈值但含文本，人工复核)")
        print(f"[extract]  page {i}: slide-{len(slides):02d}.png (dark={r:.1%})"
              + (f" ⚠ {flags[0]}" if flags else ""))

    shutil.rmtree(tmpdir, ignore_errors=True)
    manifest = {
        "source": os.path.basename(src),
        "input_file": src,
        "pages_total": len(pngs),
        "video_pages_excluded": excluded,
        "slides": slides,
    }
    with open(os.path.join(work, "manifest.json"), "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)
    print(f"[extract] ✓ {os.path.basename(src)}: "
          f"{len(slides)} 张有效幻灯片 / {len(excluded)} 个视频页 → {work}")
    return manifest


def main():
    ap = argparse.ArgumentParser(description="PPT/PDF 幻灯片提取 + 视频页过滤")
    ap.add_argument("-i", dest="inputs", action="append", required=True,
                    help="输入文件（.ppt/.pptx/.pdf），可多次传")
    ap.add_argument("-o", dest="outdir", required=True, help="输出目录")
    ap.add_argument("--dark-ratio", dest="ratio", type=float, default=DEFAULT_RATIO,
                    help=f"黑像素占比阈值（默认 {DEFAULT_RATIO}）")
    ap.add_argument("--brightness-threshold", type=int, default=DEFAULT_THRESHOLD,
                    help=f"亮度阈值（默认 {DEFAULT_THRESHOLD}）")
    ap.add_argument("--dpi", type=int, default=DEFAULT_DPI,
                    help=f"渲染 dpi（默认 {DEFAULT_DPI}）")
    args = ap.parse_args()

    os.makedirs(args.outdir, exist_ok=True)
    for src in args.inputs:
        if not os.path.exists(src):
            print(f"[extract] ✗ 文件不存在: {src}")
            sys.exit(1)
        try:
            process_file(src, args.outdir, args.dpi, args.ratio,
                         args.brightness_threshold)
        except RuntimeError as e:
            print(f"[extract] ✗ {e}")
            sys.exit(1)
    print("[extract] 全部完成")


if __name__ == "__main__":
    main()
